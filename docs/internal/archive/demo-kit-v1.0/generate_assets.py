from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle


BASE = Path(__file__).resolve().parent
IMAGES = BASE.parent / "images"

NAVY = RGBColor(22, 50, 79)
NAVY_DARK = RGBColor(13, 27, 42)
BLUE = RGBColor(62, 166, 255)
WHITE = RGBColor(255, 255, 255)
TEXT = RGBColor(15, 23, 32)
MUTED = RGBColor(93, 107, 122)


SLIDES = [
    {
        "title": "Problema",
        "bullets": [
            "La ricarica condivisa in condominio richiede ordine, chiarezza e controllo.",
            "I residenti vogliono sapere quando una colonnina e disponibile.",
            "L'amministratore vuole continuita operativa e meno gestione manuale.",
            "Le informazioni frammentate creano incertezza e richieste ripetitive.",
        ],
    },
    {
        "title": "Situazione attuale (Legrand)",
        "bullets": [
            "L'infrastruttura Legrand e il punto di partenza hardware del servizio.",
            "Il dato tecnico esiste, ma non basta da solo a gestire l'esperienza condominiale.",
            "Serve una lettura piu chiara per residenti e amministratori.",
            "Condo Charge valorizza l'impianto esistente con una piattaforma di gestione.",
        ],
    },
    {
        "title": "Condo Charge",
        "subtitle": "La piattaforma intelligente per la gestione della ricarica condominiale.",
        "bullets": [
            "Monitoraggio in tempo reale",
            "Esperienza residente mobile-first",
            "Dashboard amministratore",
            "Health monitoring e comunicazione integrata",
        ],
    },
    {
        "title": "Resident experience",
        "bullets": [
            "Vista semplice e immediata delle colonnine.",
            "Storico ricariche consultabile.",
            "Notifiche utili e non invasive.",
            "Interfaccia orientata alla disponibilita del servizio.",
        ],
        "image": IMAGES / "03-resident-dashboard.png",
    },
    {
        "title": "Admin dashboard",
        "bullets": [
            "Visione unificata dello stato del condominio.",
            "Controllo operativo e lettura immediata del servizio.",
            "Supporto alla gestione RFID e ai flussi amministrativi.",
            "Strumento pensato per decisioni rapide e chiarezza gestionale.",
        ],
        "image": IMAGES / "02-admin-dashboard.png",
    },
    {
        "title": "Push Notifications",
        "bullets": [
            "Comunicazione immediata ai residenti sugli eventi rilevanti.",
            "Meno verifiche manuali e piu tempestivita.",
            "Maggiore percezione di affidabilita del servizio.",
        ],
    },
    {
        "title": "Telegram integration",
        "bullets": [
            "Un canale familiare, rapido e ad alta apertura.",
            "Aggiornamenti di servizio accessibili anche fuori dalla piattaforma.",
            "Una comunicazione piu vicina alle abitudini quotidiane dei residenti.",
        ],
    },
    {
        "title": "Health Dashboard",
        "bullets": [
            "Monitoraggio di heartbeat, polling e salute dell'agente.",
            "Individuazione rapida di anomalie operative.",
            "Più controllo sul servizio, non solo sulla schermata.",
        ],
    },
    {
        "title": "Architecture",
        "bullets": [
            "Colonnine Legrand",
            "Agente locale Windows",
            "Backend Condo Charge",
            "PWA per residenti e amministratori",
        ],
    },
    {
        "title": "Roadmap",
        "bullets": [
            "Smart Queue",
            "Prenotazione intelligente",
            "Timer di disponibilita",
            "Statistiche avanzate",
        ],
    },
    {
        "title": "Why Condo Charge",
        "bullets": [
            "Trasforma l'infrastruttura in un servizio gestibile.",
            "Aumenta chiarezza per residenti e amministratori.",
            "Migliora comunicazione, controllo e affidabilita.",
            "Parla il linguaggio operativo del condominio.",
        ],
    },
    {
        "title": "Thank you",
        "subtitle": "Richiedi una demo di Condo Charge",
        "bullets": [
            "Condo Charge",
            "La piattaforma intelligente per la gestione della ricarica condominiale.",
        ],
    },
]


def add_full_bg(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_title_block(slide, title: str, subtitle: str | None, dark: bool) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(7.5), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE if dark else NAVY
    if subtitle:
        sp = tf.add_paragraph()
        sp.text = subtitle
        sp.font.size = Pt(14)
        sp.font.color.rgb = WHITE if dark else MUTED


def add_bullets(slide, bullets: list[str], dark: bool) -> None:
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(5.3), Inches(4.8))
    tf = box.text_frame
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(10)
        p.font.size = Pt(21 if len(bullets) <= 3 else 18)
        p.font.color.rgb = WHITE if dark else TEXT


def add_side_panel(slide, title: str, dark: bool) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(8.2),
        Inches(1.1),
        Inches(4.5),
        Inches(5.5),
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(25, 56, 89) if dark else RGBColor(240, 246, 252)
    shape.line.color.rgb = RGBColor(255, 255, 255) if dark else RGBColor(220, 229, 240)
    tx = slide.shapes.add_textbox(Inches(8.55), Inches(1.55), Inches(3.8), Inches(4.8))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = "Condo Charge"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE if dark else NAVY
    p2 = tf.add_paragraph()
    p2.text = title
    p2.font.size = Pt(17)
    p2.font.color.rgb = BLUE if dark else NAVY
    p3 = tf.add_paragraph()
    p3.text = "Piattaforma intelligente per la gestione della ricarica condominiale."
    p3.font.size = Pt(14)
    p3.font.color.rgb = WHITE if dark else MUTED


def add_image(slide, image_path: Path) -> None:
    if not image_path.exists():
        return
    slide.shapes.add_picture(str(image_path), Inches(7.0), Inches(1.55), width=Inches(5.4), height=Inches(3.6))
    caption = slide.shapes.add_textbox(Inches(7.0), Inches(5.3), Inches(5.4), Inches(0.8))
    p = caption.text_frame.paragraphs[0]
    p.text = image_path.stem.replace("-", " ").title()
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = MUTED


def build_pptx() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for index, spec in enumerate(SLIDES):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        dark = index in {0, 2, 11}
        add_full_bg(slide, NAVY_DARK if dark else WHITE)
        add_title_block(slide, spec["title"], spec.get("subtitle"), dark)
        add_bullets(slide, spec["bullets"], dark)
        if spec.get("image"):
            add_image(slide, spec["image"])
        else:
            add_side_panel(slide, spec["title"], dark)

    out = BASE / "CondoCharge-Demo-Deck-v1.0.pptx"
    prs.save(out)
    return out


def build_pdf() -> Path:
    out = BASE / "CondoCharge-Brochure-v1.0.pdf"
    styles = getSampleStyleSheet()
    title = ParagraphStyle(
        "title",
        parent=styles["Title"],
        fontName="Helvetica-Bold",
        fontSize=24,
        leading=28,
        textColor=colors.HexColor("#16324f"),
        spaceAfter=8,
    )
    subtitle = ParagraphStyle(
        "subtitle",
        parent=styles["BodyText"],
        fontName="Helvetica",
        fontSize=12,
        leading=16,
        textColor=colors.HexColor("#5d6b7a"),
        spaceAfter=12,
    )
    section = ParagraphStyle(
        "section",
        parent=styles["Heading2"],
        fontName="Helvetica-Bold",
        fontSize=12,
        leading=14,
        textColor=colors.HexColor("#16324f"),
        spaceAfter=6,
        spaceBefore=10,
    )
    body = ParagraphStyle(
        "body",
        parent=styles["BodyText"],
        fontName="Helvetica",
        fontSize=10,
        leading=14,
        textColor=colors.HexColor("#0f1720"),
    )
    small = ParagraphStyle(
        "small",
        parent=body,
        fontSize=9,
        leading=12,
        textColor=colors.HexColor("#5d6b7a"),
    )

    doc = SimpleDocTemplate(
        str(out),
        pagesize=A4,
        leftMargin=16 * mm,
        rightMargin=16 * mm,
        topMargin=14 * mm,
        bottomMargin=14 * mm,
    )
    elements = [
        Paragraph("Condo Charge", title),
        Paragraph(
            "La piattaforma intelligente per la gestione delle colonnine di ricarica nei condomini.",
            subtitle,
        ),
        Paragraph("Perche Condo Charge", section),
        Paragraph(
            "Condo Charge aiuta l'amministratore a trasformare la ricarica condominiale in un servizio ordinato, "
            "monitorabile e facilmente comunicabile ai residenti.",
            body,
        ),
        Spacer(1, 6),
    ]

    feature_rows = [
        ["Stato colonnine in tempo reale", "Notifiche Push", "Telegram"],
        ["Storico ricariche", "Gestione RFID", "Dashboard amministratore"],
        ["Health Monitoring", "Roadmap chiara", "Richiedi una demo"],
    ]
    table = Table(feature_rows, colWidths=[56 * mm, 56 * mm, 56 * mm])
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor("#eef4fb")),
                ("TEXTCOLOR", (0, 0), (-1, -1), colors.HexColor("#16324f")),
                ("FONTNAME", (0, 0), (-1, -1), "Helvetica-Bold"),
                ("FONTSIZE", (0, 0), (-1, -1), 9),
                ("LEADING", (0, 0), (-1, -1), 11),
                ("BOX", (0, 0), (-1, -1), 0.5, colors.HexColor("#d9e3ef")),
                ("INNERGRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#d9e3ef")),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                ("TOPPADDING", (0, 0), (-1, -1), 8),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
            ]
        )
    )
    elements.append(Paragraph("Cosa offre", section))
    elements.append(table)
    elements.extend(
        [
            Paragraph("Valore per l'amministratore", section),
            Paragraph(
                "Una vista chiara dello stato del servizio, maggiore continuita operativa, comunicazione piu efficace "
                "verso i residenti e una piattaforma che valorizza l'infrastruttura esistente.",
                body,
            ),
            Paragraph("Coming Soon", section),
            Paragraph(
                "Smart Queue, Prenotazione intelligente, Timer di disponibilita, Statistiche avanzate.",
                body,
            ),
            Paragraph("Posizionamento", section),
            Paragraph(
                "Condo Charge non e un'app per vedere le colonnine.<br/>"
                "E una piattaforma intelligente per la gestione della ricarica condominiale.",
                body,
            ),
            Spacer(1, 10),
            Paragraph("Richiedi una demo di Condo Charge", section),
            Paragraph(
                "Materiale demo preparato per amministratori di condominio.",
                small,
            ),
        ]
    )
    doc.build(elements)
    return out


def main() -> None:
    pptx_path = build_pptx()
    pdf_path = build_pdf()
    print(f"Generated: {pptx_path.name}")
    print(f"Generated: {pdf_path.name}")


if __name__ == "__main__":
    main()
